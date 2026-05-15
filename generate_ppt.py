"""
Generates AI_Food_Genius_Presentation.pptx — 10 slides, purple theme.
Run: python generate_ppt.py
"""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
P1  = RGBColor(102, 126, 234)   # #667eea  — main purple-blue
P2  = RGBColor(118,  75, 162)   # #764ba2  — deep purple
P3  = RGBColor( 80,  55, 140)   # darker purple accent
WH  = RGBColor(255, 255, 255)
DK  = RGBColor( 44,  62,  80)   # #2c3e50  — dark text
GR  = RGBColor( 40, 167,  69)   # green for success / ensemble
LG  = RGBColor(248, 249, 250)   # light grey background
SG  = RGBColor(108, 117, 125)   # secondary grey
LV  = RGBColor(220, 215, 255)   # lavender for sub-text on dark bg

SW = Inches(13.33)
SH = Inches(7.50)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ── Primitive helpers ───────────────────────────────────────────────────────

def box(slide, x, y, w, h, fill=None, border=None, bw=0):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if border:
        sp.line.color.rgb = border
        sp.line.width = Pt(bw) if bw else Pt(1)
    else:
        sp.line.fill.background()
    return sp

def tb(slide, text, x, y, w, h,
       sz=16, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, wrap=True):
    color = color or DK
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(sz)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return t

def fig_to_slide(slide, fig, x, y, w, h):
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    buf.seek(0)
    slide.shapes.add_picture(buf, x, y, w, h)
    plt.close(fig)

def header(slide, title, sub=None):
    box(slide, 0, 0, SW, Inches(1.15), fill=P1)
    tb(slide, title,
       Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.65),
       sz=30, bold=True, color=WH)
    if sub:
        tb(slide, sub,
           Inches(0.4), Inches(0.70), Inches(12.5), Inches(0.38),
           sz=14, italic=True, color=LV)

def bg(slide):
    box(slide, 0, Inches(1.15), SW, SH - Inches(1.15), fill=LG)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
box(sl, 0, 0, SW, SH, fill=P2)
box(sl, 0, Inches(2.4), SW, Inches(2.8), fill=P3)
box(sl, 0, Inches(6.9), SW, Inches(0.6), fill=RGBColor(70, 40, 110))

tb(sl, 'AI Food Genius',
   Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.1),
   sz=56, bold=True, color=WH, align=PP_ALIGN.CENTER)

tb(sl, 'A Personalized AI-Powered Food Recommendation System',
   Inches(1), Inches(2.65), Inches(11.3), Inches(0.7),
   sz=22, color=LV, align=PP_ALIGN.CENTER)

tb(sl, 'Python  ·  Streamlit  ·  5-Model ML Ensemble  ·  Collaborative Filtering  ·  SQLite',
   Inches(1), Inches(3.5), Inches(11.3), Inches(0.5),
   sz=15, italic=True, color=RGBColor(190, 185, 230), align=PP_ALIGN.CENTER)

tb(sl, 'Random Forest  ·  XGBoost  ·  LightGBM  ·  CatBoost  ·  Neural Network',
   Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.42),
   sz=13, color=RGBColor(180, 170, 215), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS THE APP + FEATURES
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, 'What Is AI Food Genius?')
bg(sl)

tb(sl,
   'AI Food Genius recommends the most suitable recipes for a user based on their '
   'age, health condition, dietary preference, cuisine, and taste profile.\n\n'
   'It runs entirely on a local machine — no paid APIs, no cloud services required.',
   Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.4), sz=17, color=DK)

pages = [
    ('Home',            'Dashboard with platform stats\nand trending recipes'),
    ('AI Recommendations', 'Enter your profile and get\npersonalized recipe suggestions'),
    ('Analytics',       'Model accuracy charts,\ncalorie & cuisine breakdowns'),
    ('Profile',         'Save preferences permanently;\nview your rating history'),
    ('Favourites',      'All recipes you saved,\nin one place'),
]
bx = Inches(0.4)
for name, desc in pages:
    box(sl, bx, Inches(3.0), Inches(2.35), Inches(2.1),
        fill=WH, border=P1, bw=1.5)
    tb(sl, name, bx + Inches(0.1), Inches(3.1), Inches(2.15), Inches(0.42),
       sz=13, bold=True, color=P2, align=PP_ALIGN.CENTER)
    tb(sl, desc, bx + Inches(0.1), Inches(3.6), Inches(2.15), Inches(0.85),
       sz=11, color=DK, align=PP_ALIGN.CENTER)
    bx += Inches(2.55)

box(sl, Inches(0.3), Inches(5.5), Inches(12.7), Inches(0.55), fill=RGBColor(232, 250, 235))
tb(sl, 'Users register and log in — profile, favourites, and ratings all persist across sessions (SQLite database).',
   Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.45),
   sz=14, bold=True, color=GR)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE DATASET
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, 'The Dataset', 'Structured synthetic dataset built specifically for this application')
bg(sl)

# 4 stat boxes
stats = [('250', 'Unique Recipes'), ('200', 'User Profiles'),
         ('3,000', 'Interactions'),  ('26',  'Columns / Record')]
bx = Inches(0.5)
for val, label in stats:
    box(sl, bx, Inches(1.35), Inches(2.8), Inches(1.75), fill=WH, border=P1, bw=2)
    tb(sl, val,   bx, Inches(1.55), Inches(2.8), Inches(0.75),
       sz=38, bold=True, color=P1, align=PP_ALIGN.CENTER)
    tb(sl, label, bx, Inches(2.38), Inches(2.8), Inches(0.5),
       sz=13, color=DK, align=PP_ALIGN.CENTER)
    bx += Inches(3.1)

rows_data = [
    ('Cuisines',           'Indian  ·  Italian  ·  Asian  ·  Mediterranean  ·  Mexican  ·  American'),
    ('Diet Types',         'Vegetarian  ·  Non-Vegetarian  ·  Vegan  ·  Keto  ·  Paleo'),
    ('Taste Profiles',     'Spicy  ·  Sweet  ·  Mild  ·  Savory  ·  Tangy'),
    ('Health Conditions',  'None  ·  Diabetes  ·  Iron Deficiency  ·  Vitamin D Deficiency  ·  High Blood Pressure'),
]
ty = Inches(3.5)
for i, (cat, vals) in enumerate(rows_data):
    bg_col = WH if i % 2 == 0 else RGBColor(238, 235, 252)
    box(sl, Inches(0.4), ty, Inches(12.5), Inches(0.48), fill=bg_col, border=RGBColor(220,215,240), bw=0.5)
    tb(sl, cat + ':', Inches(0.55), ty + Inches(0.07), Inches(2.5), Inches(0.38),
       sz=13, bold=True, color=P2)
    tb(sl, vals,      Inches(3.2),  ty + Inches(0.07), Inches(9.5), Inches(0.38),
       sz=13, color=DK)
    ty += Inches(0.5)

tb(sl, 'Every cuisine–diet combination has 4–15 recipes. '
        'Nutritional values are realistic by diet type (e.g. Keto = high fat & low carb by design).',
   Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
   sz=13, italic=True, color=SG)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RECOMMENDATION PIPELINE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, 'How Recommendations Work', '4-layer pipeline — runs every time you click Get Recommendations')
bg(sl)

steps = [
    (P1,                         '1', 'Filter',        'Diet · Cuisine\nCooking Time'),
    (P2,                         '2', 'ML Score',      '5-Model\nWeighted Ensemble'),
    (P3,                         '3', 'Personalize',   'Taste · Health\nNutrition'),
    (RGBColor(60, 40, 130),      '4', 'Collab Filter', 'Similar-User\nRatings'),
    (GR,                         '✓', 'Top 10',        'Best Recipes\nShown to User'),
]
bx = Inches(0.4)
for i, (col, num, title, sub) in enumerate(steps):
    box(sl, bx, Inches(1.9), Inches(2.2), Inches(2.05), fill=col)
    tb(sl, num,   bx, Inches(2.0),  Inches(2.2), Inches(0.55),
       sz=28, bold=True, color=WH, align=PP_ALIGN.CENTER)
    tb(sl, title, bx, Inches(2.6),  Inches(2.2), Inches(0.5),
       sz=15, bold=True, color=WH, align=PP_ALIGN.CENTER)
    tb(sl, sub,   bx, Inches(3.15), Inches(2.2), Inches(0.65),
       sz=11, color=LV, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        tb(sl, '→', bx + Inches(2.2), Inches(2.7), Inches(0.38), Inches(0.45),
           sz=24, bold=True, color=P2, align=PP_ALIGN.CENTER)
    bx += Inches(2.6)

desc_lines = [
    '① Only recipes matching the selected diet type and cuisine enter — a Non-Vegetarian filter will never show chicken-free results.',
    '② All 5 ML models score each recipe individually. Each model\'s vote is weighted by its measured accuracy.',
    '③ Taste match (+25%), health-condition nutrition bonus (+20%), and nutrient quality (+15%) are added on top of the ML score.',
    '④ After rating recipes, similar-user suggestions appear automatically — same approach used by Netflix and Spotify.',
]
ty = Inches(4.45)
for line in desc_lines:
    tb(sl, line, Inches(0.5), ty, Inches(12.5), Inches(0.44), sz=13, color=DK)
    ty += Inches(0.5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 5-MODEL ENSEMBLE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, 'Layer 2 — The 5-Model AI Ensemble', 'Each model scores every recipe for the user; results are accuracy-weighted and combined')
bg(sl)

models = [
    (P1,                    'Random Forest',  'scikit-learn',
     '100 trees · max depth 10',
     'An ensemble of 100 decision trees each trained on a random data subset. They vote together — robust and noise-resistant.'),
    (P2,                    'XGBoost',        'xgboost',
     '100 trees · depth 6 · logloss',
     'Gradient-boosted trees built sequentially — each new tree corrects the previous one\'s errors. Strong on tabular data.'),
    (P3,                    'Neural Network', 'scikit-learn MLP',
     '64 → 32 neurons · 300 iterations',
     'A three-layer neural network that learns non-linear relationships between features that tree models may miss.'),
    (RGBColor(60,100,200),  'LightGBM',       'Microsoft',
     '200 trees · lr 0.05 · 31 leaves',
     'Leaf-wise tree growth (faster & more accurate than level-wise XGBoost). Handles large feature sets well.'),
    (RGBColor(100,60,180),  'CatBoost',       'Yandex',
     '200 iterations · depth 6 · lr 0.05',
     'Handles categorical features (cuisine, diet type) natively — no manual encoding needed, less information loss.'),
]
ty = Inches(1.35)
for col, name, dev, params, desc in models:
    box(sl, Inches(0.35), ty, Inches(2.7), Inches(0.92), fill=col)
    tb(sl, name, Inches(0.45), ty + Inches(0.05), Inches(2.55), Inches(0.4),
       sz=14, bold=True, color=WH)
    tb(sl, f'by {dev}', Inches(0.45), ty + Inches(0.48), Inches(2.55), Inches(0.35),
       sz=11, italic=True, color=LV)
    tb(sl, params, Inches(3.15), ty + Inches(0.04), Inches(3.0), Inches(0.38),
       sz=12, bold=True, color=P2)
    tb(sl, desc, Inches(3.15), ty + Inches(0.44), Inches(9.7), Inches(0.48),
       sz=12, color=DK)
    ty += Inches(1.04)

box(sl, Inches(0.35), Inches(6.78), Inches(12.6), Inches(0.52), fill=GR)
tb(sl, 'Weighted Ensemble — Final score = Σ (model_accuracy / total_accuracy) × model_probability',
   Inches(0.5), Inches(6.82), Inches(12.3), Inches(0.42),
   sz=14, bold=True, color=WH, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SCORING BREAKDOWN (PIE CHART)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, 'Layer 3 — Multi-Component Personalized Scoring', 'ML probability is only 40% — three more signals are combined on top')
bg(sl)

# Pie chart
fig, ax = plt.subplots(figsize=(5.8, 5.0))
sizes   = [40, 25, 20, 15]
labels  = ['ML Ensemble\nProbability\n40%', 'Taste\nMatch\n25%',
           'Health\nCondition\nBonus\n20%', 'Nutrient\nQuality\n15%']
colors  = ['#667eea', '#764ba2', '#28a745', '#f39c12']
explode = (0.05, 0.05, 0.05, 0.05)
wedges, texts = ax.pie(sizes, explode=explode, labels=labels, colors=colors,
                       startangle=90, textprops={'fontsize': 11, 'fontweight': 'bold'})
ax.set_title('Final Score Composition', fontsize=13, fontweight='bold', pad=14)
fig.patch.set_facecolor('white')
fig_to_slide(sl, fig, Inches(0.2), Inches(1.3), Inches(6.3), Inches(5.6))

# Details panel
details = [
    (P1,  '40%  ML Ensemble',
     'Probability output from 5 trained models\nlearned from 3,000 user–recipe interactions.'),
    (P2,  '25%  Taste Match',
     'Recipes whose taste tag (Spicy / Sweet / Mild /\nSavory / Tangy) matches yours get +0.25 bonus.'),
    (GR,  '20%  Health Condition Bonus',
     'Diabetes  →  low carb + high fiber\nIron Deficiency  →  high iron foods\nHigh Blood Pressure  →  low fat + high fiber'),
    (RGBColor(200, 130, 20), '15%  Nutrient Quality Score',
     'A recipe-level nutritional balance score.\nHigher quality recipes ranked above lower ones.'),
]
ty = Inches(1.4)
for col, title_d, body_d in details:
    box(sl, Inches(6.55), ty, Inches(0.1), Inches(0.9), fill=col)
    tb(sl, title_d, Inches(6.75), ty + Inches(0.02), Inches(6.2), Inches(0.38),
       sz=13, bold=True, color=col)
    tb(sl, body_d, Inches(6.75), ty + Inches(0.38), Inches(6.2), Inches(0.58),
       sz=11, color=DK)
    ty += Inches(1.2)

tb(sl, 'Scores are normalized within the result set — the best and worst recipes are always clearly separated.',
   Inches(6.55), Inches(6.75), Inches(6.5), Inches(0.5),
   sz=12, italic=True, color=SG)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ACCURACY BAR CHART
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, 'Model Accuracy', 'Evaluated on 20% held-out test set (600 records) — User_Rating excluded to prevent data leakage')
bg(sl)

fig, ax = plt.subplots(figsize=(7.2, 4.8))
m_names = ['Random\nForest', 'CatBoost', 'XGBoost', 'LightGBM', 'Neural\nNet', 'Weighted\nEnsemble']
accs    = [62.8, 59.3, 57.3, 56.3, 53.5, 63.0]
bcolors = ['#667eea', '#764ba2', '#5563d4', '#8050b0', '#9b59b6', '#28a745']
bars = ax.bar(m_names, accs, color=bcolors, edgecolor='white', linewidth=1.5, width=0.55)
ax.set_ylim(44, 68)
ax.set_ylabel('Accuracy (%)', fontsize=12)
ax.set_title('Classification Accuracy — will user recommend this recipe?', fontsize=12, fontweight='bold')
ax.axhline(y=63, color='#28a745', linestyle='--', linewidth=1.8, alpha=0.8, label='Ensemble ~63%')
ax.legend(fontsize=11)
for bar, acc in zip(bars, accs):
    ax.text(bar.get_x() + bar.get_width() / 2., bar.get_height() + 0.3,
            f'{acc}%', ha='center', va='bottom', fontweight='bold', fontsize=11)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
fig.patch.set_facecolor('white')
fig.tight_layout()
fig_to_slide(sl, fig, Inches(0.3), Inches(1.35), Inches(7.5), Inches(5.6))

# Notes panel
notes = [
    ('What we measured',
     'Binary classification: will this user\nrecommend this recipe? (yes/no)'),
    ('Test setup',
     '600 records held out — never seen\nduring model training'),
    ('Best single model',
     'Random Forest at 62.8%'),
    ('Ensemble result',
     '~63% — beats all individual models\nthrough accuracy-weighted voting'),
    ('Data leakage prevented',
     'User_Rating excluded from features\n(it directly predicts the target)'),
]
ty = Inches(1.4)
tb(sl, 'Key Notes', Inches(8.1), ty - Inches(0.1), Inches(4.9), Inches(0.42),
   sz=16, bold=True, color=P2)
ty += Inches(0.45)
for title_n, body_n in notes:
    tb(sl, title_n, Inches(8.1), ty, Inches(4.9), Inches(0.35),
       sz=13, bold=True, color=P2)
    tb(sl, body_n, Inches(8.1), ty + Inches(0.32), Inches(4.9), Inches(0.55),
       sz=12, color=DK)
    ty += Inches(1.0)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — VS RESEARCH PAPERS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, 'Our Results vs Research Papers', 'Honest context — why a direct comparison is not valid')
bg(sl)

# Comparison table
col_ws = [Inches(2.8), Inches(4.0), Inches(5.1)]
headers_t = ['Factor', 'Our System', 'Research Papers']
rows_t = [
    ['Dataset',          '3,000 interactions (synthetic)', '1 million+ interactions (real)'],
    ['Recipes',          '250 unique (generated)',          '100,000 – 1 million (real)'],
    ['Users',            '200 (generated profiles)',        'Millions of real users'],
    ['Evaluation metric','Classification accuracy',         'NDCG@10, Precision@K, MAP'],
    ['Reported accuracy','~63%',                            '75 – 90%'],
    ['Data source',      'Synthetically generated',         'Years of real user behaviour'],
]
tx = Inches(0.35)
ty = Inches(1.4)
for j, h in enumerate(headers_t):
    box(sl, tx + sum(col_ws[:j]), ty, col_ws[j], Inches(0.48), fill=P1)
    tb(sl, h, tx + sum(col_ws[:j]) + Inches(0.1), ty + Inches(0.07),
       col_ws[j] - Inches(0.1), Inches(0.38), sz=13, bold=True, color=WH)

for i, row in enumerate(rows_t):
    ty2 = ty + Inches(0.48 * (i + 1))
    bg_row = WH if i % 2 == 0 else RGBColor(238, 235, 252)
    for j, cell in enumerate(row):
        box(sl, tx + sum(col_ws[:j]), ty2, col_ws[j], Inches(0.46),
            fill=bg_row, border=RGBColor(220, 215, 240), bw=0.5)
        cell_col = GR if (j == 2 and i == 4) else (RGBColor(190, 50, 50) if (j == 1 and i == 4) else DK)
        tb(sl, cell,
           tx + sum(col_ws[:j]) + Inches(0.1), ty2 + Inches(0.06),
           col_ws[j] - Inches(0.1), Inches(0.38),
           sz=12, color=cell_col, bold=(j == 0))

box(sl, Inches(0.35), Inches(5.28), Inches(12.6), Inches(0.06), fill=P2)

tb(sl, 'What to say when asked:',
   Inches(0.35), Inches(5.45), Inches(12.5), Inches(0.4),
   sz=15, bold=True, color=P2)

quote = ('"Our system achieves ~63% accuracy on a 3,000-record synthetic dataset using a 5-model weighted ensemble. '
         'Research papers report 75–90% on million-record real-world datasets using ranking metrics like NDCG@10 — '
         'a direct comparison is not meaningful at this scale. '
         'The architecture is production-grade; connecting to real data would close the gap significantly."')
box(sl, Inches(0.35), Inches(5.9), Inches(12.6), Inches(1.3), fill=RGBColor(240, 238, 255))
tb(sl, quote, Inches(0.5), Inches(5.96), Inches(12.3), Inches(1.2),
   sz=13, italic=True, color=DK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DATA STORAGE & SECURITY
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, 'Data Storage & Security')
bg(sl)

# Two columns
col_items = [
    ('CSV Files  (final_dataset.csv)', P1, [
        '250 recipes × 200 users × 3,000 interactions',
        '26 columns: profile + recipe nutrition + scores',
        'Generated by h1.py — re-runnable anytime',
        'What the 5 ML models train on at startup',
    ]),
    ('SQLite Database  (food_genius.db)', P2, [
        'users — username, SHA-256 hashed password, full profile',
        'favourites — saved recipe names per user + timestamp',
        'ratings — 1–5 star rating per user-recipe pair + timestamp',
        'Passwords NEVER stored in plain text',
    ]),
]
bx = Inches(0.4)
for title_c, col, items in col_items:
    box(sl, bx, Inches(1.35), Inches(6.0), Inches(0.48), fill=col)
    tb(sl, title_c, bx + Inches(0.1), Inches(1.38), Inches(5.8), Inches(0.4),
       sz=14, bold=True, color=WH)
    ty = Inches(1.95)
    for item in items:
        box(sl, bx, ty, Inches(6.0), Inches(0.6),
            fill=WH if items.index(item) % 2 == 0 else RGBColor(238, 235, 252),
            border=RGBColor(220, 215, 240), bw=0.5)
        tb(sl, '●  ' + item, bx + Inches(0.1), ty + Inches(0.1),
           Inches(5.8), Inches(0.5), sz=13, color=DK)
        ty += Inches(0.62)
    bx += Inches(6.6)

tb(sl, 'How the system learns from ratings:',
   Inches(0.4), Inches(5.35), Inches(12.5), Inches(0.4),
   sz=15, bold=True, color=P2)
learn = ('Every time a logged-in user rates a recipe (1–5 stars), that rating is saved to the database. '
         'These ratings feed directly into the collaborative filtering engine — the more you rate, the better the '
         '"users like you" recommendations become.\n'
         'The ML models train once at startup and are cached. They benefit from new ratings on the next session restart — '
         'the same way production systems retrain on a daily/weekly schedule.')
tb(sl, learn, Inches(0.4), Inches(5.82), Inches(12.5), Inches(1.35), sz=13, color=DK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TECH STACK + SUMMARY
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header(sl, 'Technology Stack & Summary')
bg(sl)

# Tech stack table
stack = [
    ('Web Framework',     'Streamlit'),
    ('Data Handling',     'Pandas, NumPy'),
    ('Machine Learning',  'scikit-learn · XGBoost · LightGBM · CatBoost'),
    ('Similarity Search', 'scikit-learn — cosine similarity'),
    ('Visualisation',     'Plotly'),
    ('Database',          'SQLite  (Python built-in — no install needed)'),
    ('Language',          'Python 3'),
]
tb(sl, 'Technology Stack', Inches(0.4), Inches(1.3), Inches(6.0), Inches(0.42),
   sz=17, bold=True, color=P2)
ty = Inches(1.8)
for i, (layer, tech) in enumerate(stack):
    bg_c = WH if i % 2 == 0 else RGBColor(238, 235, 252)
    box(sl, Inches(0.4), ty, Inches(6.0), Inches(0.44),
        fill=bg_c, border=RGBColor(220, 215, 240), bw=0.5)
    tb(sl, layer, Inches(0.55), ty + Inches(0.06), Inches(2.0), Inches(0.35),
       sz=12, bold=True, color=P2)
    tb(sl, tech,  Inches(2.6),  ty + Inches(0.06), Inches(3.65), Inches(0.35),
       sz=12, color=DK)
    ty += Inches(0.46)

# Summary bullets
tb(sl, 'Summary', Inches(7.1), Inches(1.3), Inches(6.0), Inches(0.42),
   sz=17, bold=True, color=P2)
summary = [
    '5-model weighted AI ensemble (RF, XGBoost, LightGBM, CatBoost, MLP)',
    '4-layer recommendation pipeline: Filter → Score → Personalize → Collab',
    '250 real-named recipes across 6 cuisines and 5 diet types',
    'Health-aware scoring: Diabetes, Iron Deficiency, High BP, Vitamin D',
    'Collaborative filtering — Netflix-style "users like you" recommendations',
    'Persistent login, favourites, and ratings via SQLite',
    'Passwords hashed with SHA-256 — never stored in plain text',
    '~63% accuracy on held-out synthetic test set (600 records)',
    'Architecture scales to research-level performance with real data',
]
ty = Inches(1.8)
for pt in summary:
    tb(sl, '●  ' + pt, Inches(7.1), ty, Inches(5.9), Inches(0.42), sz=12, color=DK)
    ty += Inches(0.5)

# Footer
box(sl, 0, Inches(7.1), SW, Inches(0.4), fill=P2)
tb(sl, 'AI Food Genius  |  Python · Streamlit · 5-Model Ensemble · Collaborative Filtering · SQLite',
   Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.35),
   sz=12, color=WH, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
prs.save('AI_Food_Genius_Presentation.pptx')
print(f'Saved: AI_Food_Genius_Presentation.pptx  ({len(prs.slides)} slides)')
